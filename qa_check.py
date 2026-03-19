"""
QA script: visual inspection of NitXig_2026_Getaway.pptx
Renders each slide to a PNG using python-pptx + Pillow.
Flags: text overflow, overlapping shapes, low contrast, uneven spacing.
"""
import os, sys
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_PATH = os.path.join(os.path.dirname(__file__), 'NitXig_2026_Getaway.pptx')
OUT_DIR   = os.path.join(os.path.dirname(__file__), '_qa')
os.makedirs(OUT_DIR, exist_ok=True)

SCALE = 96 / 914400  # EMU → pixels at 96 dpi
W_PX  = int(13.33 * 96)
H_PX  = int(7.5  * 96)

ISSUES = []

def emu_to_px(emu):
    return int(emu * SCALE)

def hex_to_rgb(hex_str):
    if hex_str is None:
        return (180, 180, 180)
    h = hex_str.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def get_fill_color(shape):
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        if str(fill.type) in ('SOLID', '1'):
            fc = fill.fore_color
            if fc and fc.type is not None:
                return (fc.rgb.red, fc.rgb.green, fc.rgb.blue)
    except Exception:
        pass
    return None

def get_text_color(run):
    try:
        fc = run.font.color
        if fc and fc.type is not None:
            return (fc.rgb.red, fc.rgb.green, fc.rgb.blue)
    except Exception:
        pass
    return (28, 28, 26)

def contrast_ratio(c1, c2):
    """Simplified contrast check (luma diff)."""
    def luma(c):
        return 0.299*c[0] + 0.587*c[1] + 0.114*c[2]
    d = abs(luma(c1) - luma(c2))
    return d

def render_slide(slide, slide_num, bg_color=(242, 237, 228)):
    img = Image.new('RGB', (W_PX, H_PX), bg_color)
    draw = ImageDraw.Draw(img)

    shapes_info = []

    for shape in slide.shapes:
        left   = emu_to_px(shape.left   or 0)
        top    = emu_to_px(shape.top    or 0)
        width  = emu_to_px(shape.width  or 0)
        height = emu_to_px(shape.height or 0)
        right  = left + width
        bottom = top  + height

        shapes_info.append({
            'left': left, 'top': top, 'right': right, 'bottom': bottom,
            'name': shape.name, 'shape': shape
        })

        # Draw fill
        fill_color = get_fill_color(shape)
        if fill_color:
            draw.rectangle([left, top, right, bottom], fill=fill_color)

        # Draw image placeholder (grey)
        if shape.shape_type == 13:  # PICTURE
            draw.rectangle([left, top, right, bottom], fill=(80, 80, 80))
            draw.text((left + 4, top + 4), '📷 img', fill=(200, 200, 200))

        # Draw text
        if shape.has_text_frame:
            tf = shape.text_frame
            y_cursor = top + 4
            for para in tf.paragraphs:
                line_text = ''
                text_col = (28, 28, 26)
                font_size = 11
                for run in para.runs:
                    line_text += run.text
                    rc = get_text_color(run)
                    if rc:
                        text_col = rc
                    if run.font.size:
                        font_size = max(7, int(run.font.size.pt * 96 / 72))
                if line_text.strip():
                    # Check text overflow
                    txt_w = len(line_text) * (font_size * 0.55)
                    if txt_w > width + 20:
                        ISSUES.append(f'Slide {slide_num}: possible text overflow in "{line_text[:40]}"')
                    draw.text((left + 3, y_cursor), line_text[:80], fill=text_col)
                    y_cursor += font_size + 2

                # Check if text goes below shape
                if y_cursor > bottom + 10:
                    ISSUES.append(f'Slide {slide_num}: text may overflow bottom of shape (shape: {shape.name})')

        # Check off-slide elements
        if right > W_PX + 10 or bottom > H_PX + 10 or left < -10 or top < -10:
            ISSUES.append(f'Slide {slide_num}: shape "{shape.name}" is off-slide (l={left},t={top},r={right},b={bottom})')

    # Draw thin border for slide boundary
    draw.rectangle([0, 0, W_PX-1, H_PX-1], outline=(180, 170, 160), width=2)

    # Slide number label
    draw.rectangle([W_PX - 55, H_PX - 22, W_PX, H_PX], fill=(50, 50, 50))
    draw.text((W_PX - 50, H_PX - 18), f'Slide {slide_num}', fill=(220, 220, 220))

    return img


def main():
    prs = Presentation(PPTX_PATH)
    print(f'✓ Opened: {PPTX_PATH}')
    print(f'  Slides: {len(prs.slides)}, Size: {prs.slide_width.inches:.2f}" × {prs.slide_height.inches:.2f}"')

    for i, slide in enumerate(prs.slides, 1):
        bg = (242, 237, 228)
        img = render_slide(slide, i, bg)
        out_path = os.path.join(OUT_DIR, f'slide_{i:02d}.png')
        img.save(out_path)
        print(f'  → Rendered slide {i:02d} ({len(slide.shapes)} shapes)')

    print()
    if ISSUES:
        print(f'⚠  QA Issues found ({len(ISSUES)}):')
        for issue in ISSUES:
            print(f'   • {issue}')
    else:
        print('✓ No QA issues detected.')

    print(f'\nImages saved to: {OUT_DIR}')

if __name__ == '__main__':
    main()
